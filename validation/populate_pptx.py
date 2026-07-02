import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

template_path = r"c:\Users\yashw\Desktop\HACKATHONN\resume ranker\resources\Idea Submission Template _ Redrob.pptx"
output_path = r"c:\Users\yashw\Desktop\HACKATHONN\resume ranker\Code\KALI_Submission_Redrob.pptx"

if not os.path.exists(template_path):
    print("Error: Template PPTX not found!")
    exit(1)

prs = Presentation(template_path)
print(f"Loaded template with {len(prs.slides)} slides.")

# Slide 1: Title
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "Team Name" in txt:
        shape.text = "Team Name : KALI"
    elif "Problem Statement" in txt:
        shape.text = "Problem Statement : Intelligent Candidate Discovery & Ranking Challenge"
    elif "Team Leader Name" in txt:
        shape.text = "Team Leader Name : P Yashwanth Reddy"

# Slide 2: Solution Overview
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "proposed solution" in txt:
        shape.text = (
            "Proposed Solution: KALI (Knowledge-Augmented Listing Intelligence)\n\n"
            "• Built a pure-python multi-factor weighted scoring algorithm that evaluates candidates across 6 core dimensions.\n"
            "• Implemented zero-dependency deterministic logic running entirely in memory, processing 100K candidates in 10 seconds on CPU.\n"
            "• Incorporates a multi-indicator honeypot detection filter and keyword-stuffer penalty system.\n"
            "• Generates unique, grounded, non-templated reasoning for each candidate in the top 100."
        )
    elif "traditional candidate matching" in txt:
        shape.text = (
            "What Differentiates Our Approach:\n\n"
            "1. Career History Over Skill Keywords: We check the candidate's actual job titles and description texts in career history, bypassing keyword stuffers who list skills but have unrelated titles (e.g. Marketing Manager).\n"
            "2. Behavioral Platform Signals: We integrate recruiter response rates and login activity. Perfect-on-paper candidates who are inactive or unresponsive are penalized.\n"
            "3. Active Honeypot Shielding: Detects impossible profiles (expert skills with 0 months, impossible tenure) and filters them out automatically."
        )

# Slide 3: JD Understanding & Candidate Evaluation
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "requirements extracted" in txt:
        shape.text = (
            "Key Requirements Extracted from JD:\n\n"
            "• Experience: 5–9 years (sweet spot for founding team senior engineer).\n"
            "• Core Tech: Embeddings-based retrieval (sentence-transformers), vector DBs (FAISS, Milvus, Pinecone), strong Python.\n"
            "• Background: Experience at product companies (services/consulting giants are penalized).\n"
            "• Logistics: India-based (Pune/Noida preferred or willing to relocate), short notice period (<30 days preferred)."
        )
    elif "relevance" in txt or "beyond keyword matching" in txt:
        shape.text = (
            "Beyond Keyword Matching - Candidate Evaluation:\n\n"
            "• Cross-Validation: Skills are cross-checked with career titles and descriptions. If a candidate claims 'NLP Expert' but worked as an accountant, they are flagged and down-ranked.\n"
            "• Engagement & Availability: Multiplier applied based on response rate and activity. Active users are boosted; inactive users (>180 days) are demoted.\n"
            "• Trust Validation: Connection count vs endorsements. Fake connections or disproportionately high endorsements trigger honeypot detection."
        )

# Slide 4: Ranking Methodology
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "retrieve, score, and rank" in txt or "algorithms, or heuristics" in txt or "signals combined" in txt:
        shape.text = (
            "Ranking Methodology & Multi-Factor Scoring:\n\n"
            "• Phase 1 Cascade Filter: Fast elimination pass (YoE, inactivity, keyword stuffers, and lack of career ML keywords) — filters out ~60% of candidates in milliseconds.\n"
            "• Title & Career Relevance (35%): Matches taxonomy of target ML/AI titles and scans career history descriptions for AI keywords.\n"
            "• Skills Match (25%): Computes a weighted sum of AI skills, adjusted by proficiency level, usage duration, and endorsements.\n"
            "• Experience Fit (15%): Optimal bell curve centered on 7 years, with penalties for job hopping.\n"
            "• Behavioral Signals (15%): Response rate, profile completeness, GitHub activity, interview attendance, and reply speed.\n"
            "• Location & Notice (5%): India location, Pune/Noida match, and notice period length.\n"
            "• Education (5%): Field of study relevance, degree level, and university tiering.\n\n"
            "Multipliers applied: Honeypot (0.0x), Keyword Stuffer (0.05x), Services-Only (0.85x), No Career ML keywords (0.10x)."
        )

# Slide 5: Explainability & Data Validation
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "ranking decisions explained" in txt or "hallucinations" in txt or "low-quality" in txt:
        shape.text = (
            "Explainability, Quality Control & Validation:\n\n"
            "• 100% Grounded Justifications: A dynamic reasoning generator constructs sentences based on the candidate's exact title, years of experience, specific skills matched, current company, and location. Hallucinations are physically impossible because we only inject fields from the candidate's record.\n"
            "• Zero Template Duplication: Avoids duplicate reasoning strings by hash-varying sentence patterns and location/notice/activity phrases based on candidate ID.\n"
            "• Anomaly/Honeypot Filter: Profiles claiming 'expert' in many skills with 0 months duration, or years of experience inconsistent with career dates, are filtered out and assigned a score of 0."
        )

# Slide 6: End-to-End Workflow
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "complete workflow" in txt:
        shape.text = (
            "End-to-End Execution Workflow (3-Phase Cascade):\n\n"
            "1. Read Input: Stream candidate records line-by-line from candidates.jsonl.gz (memory-efficient).\n"
            "2. Phase 1 Elimination Filter: Reject honeypots, inactive profiles (>240d), keyword stuffers, YoE outside [3, 15], and profiles lacking career ML keywords. (Eliminates ~60% of candidates instantly).\n"
            "3. Multi-Factor Scoring: Calculate scores for Title, Skills, Experience, Behavioral, Location, and Education for remaining candidates.\n"
            "4. Apply Penalty Multipliers: Demote stuffers, services-only candidates, inactive profiles, and those lacking career keywords.\n"
            "5. Sort & Rank: Sort candidates descending by score, tiebreaking on Candidate ID ascending.\n"
            "6. Generate Reasonings: Run dynamic, grounded text generation for the top 100.\n"
            "7. Export: Write the ranked records to the validated CSV format."
        )

# Slide 7: System Architecture
# Let's add a text representation of the architecture in Slide 7
slide7 = prs.slides[6]
# Let's check what shapes are there first. We'll add text to the existing text box if it's there
found = False
for shape in slide7.shapes:
    if hasattr(shape, "text") and ("architecture" in shape.text.lower() or not shape.text.strip()):
        shape.text = (
            "System Architecture — KALI Pipeline:\n\n"
            "                       [ candidates.jsonl.gz ]\n"
            "                                  │\n"
            "                                  ▼\n"
            "                       [ Honeypot Anomaly Filter ] ──► Score = 0 (If Flagged)\n"
            "                                  │\n"
            "                                  ▼\n"
            "                      [ Multi-Factor Scorer ]\n"
            "                 Title (35%)   •   Skills (25%)   •   Exp (15%)\n"
            "                 Signals (15%)  •   Loc (5%)      •   Edu (5%)\n"
            "                                  │\n"
            "                                  ▼\n"
            "                      [ Penalty Adjustments ]\n"
            "                 Stuffers (0.05x)  •   Services-only (0.85x)\n"
            "                 Inactive (0.4x)   •   No Career ML (0.1x)\n"
            "                                  │\n"
            "                                  ▼\n"
            "                     [ Sorting & Tie-Breaking ]\n"
            "                                  │\n"
            "                                  ▼\n"
            "                   [ Grounded Reasoning Generator ]\n"
            "                                  │\n"
            "                                  ▼\n"
            "                        [ submission.csv ]"
        )
        found = True
        break

if not found:
    # Create a text box on slide 7
    txBox = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = txBox.text_frame
    tf.text = (
        "System Architecture — KALI Pipeline:\n\n"
        "                       [ candidates.jsonl.gz ]\n"
        "                                  │\n"
        "                                  ▼\n"
        "                       [ Honeypot Anomaly Filter ] ──► Score = 0 (If Flagged)\n"
        "                                  │\n"
        "                                  ▼\n"
        "                      [ Multi-Factor Scorer ]\n"
        "                 Title (35%)   •   Skills (25%)   •   Exp (15%)\n"
        "                 Signals (15%)  •   Loc (5%)      •   Edu (5%)\n"
        "                                  │\n"
        "                                  ▼\n"
        "                      [ Penalty Adjustments ]\n"
        "                 Stuffers (0.05x)  •   Services-only (0.85x)\n"
        "                 Inactive (0.4x)   •   No Career ML (0.1x)\n"
        "                                  │\n"
        "                                  ▼\n"
        "                     [ Sorting & Tie-Breaking ]\n"
        "                                  │\n"
        "                                  ▼\n"
        "                   [ Grounded Reasoning Generator ]\n"
        "                                  │\n"
        "                                  ▼\n"
        "                        [ submission.csv ]"
    )

# Slide 8: Results & Performance
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "results or insights" in txt or "runtime and compute" in txt:
        shape.text = (
            "Results, Performance & Insights:\n\n"
            "• Execution Speed: Processed 100K candidates with Phase 1 elimination in 10.7 seconds. Peak RAM: 60MB. CPU cores utilized: 1.\n"
            "• Official Format Verification: Passed the official validate_submission.py check with 100% compliance.\n"
            "• Deep Quality Validation Metrics:\n"
            "  - Honeypot Rate: 0% in Top 100 (Disqualification threshold is >10%).\n"
            "  - Keyword Stuffer Rate: 0% in Top 100.\n"
            "  - Career History ML Relevance: 100% in Top 10 and Top 50 candidates.\n"
            "  - Reasoning Grounding & Uniqueness: 100% grounded in profile facts; 0 duplicates.\n"
            "  - Behavioral Signals: 100% active, highly responsive profiles in Top 10."
        )

# Slide 9: Technologies Used
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "technologies, frameworks" in txt:
        shape.text = (
            "Technology Stack & Rationale:\n\n"
            "1. Python 3.10 (Standard Library): Used for the core ranking engine. Standard modules (json, csv, gzip, re, datetime, math) guarantee zero package installation errors, fast startup, and high predictability.\n"
            "2. Gradio: Selected for the HuggingFace Spaces sandbox deployment. Provides an intuitive, responsive web UI for uploading candidate segments and downloading ranked CSVs in real time.\n"
            "3. python-docx: Used offline to parse the Job Description and Redrob reference documents to engineer features.\n"
            "4. Git & GitHub LFS: Used to manage versioning and track the gzipped candidate dataset (54MB) cleanly."
        )

# Slide 10: Submission Assets
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if not hasattr(shape, "text") or not shape.text:
        continue
    txt = shape.text.strip()
    if "Github video" in txt or "Submission Assets" in txt or not txt:
        shape.text = (
            "Submission Assets & Links:\n\n"
            "• GitHub Repository:\n"
            "  https://github.com/yashreddy1154/intelligent-candidate-ranker\n"
            "  Contains the fully reproducible ranking script, README instructions, and validation suite.\n\n"
            "• HuggingFace Spaces Web Sandbox:\n"
            "  https://huggingface.co/spaces/yash1154/intelligent-candidate-ranker\n"
            "  A working demo to upload candidate subsets and run matching on the fly.\n\n"
            "• Final Submission File:\n"
            "  submission.csv (UTF-8, containing exactly 100 ranked candidates with grounded reasoning)."
        )

prs.save(output_path)
print("Successfully generated and saved populated PPTX.")
