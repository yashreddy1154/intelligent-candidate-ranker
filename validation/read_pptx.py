import os
from pptx import Presentation

pptx_path = r"c:\Users\yashw\Desktop\HACKATHONN\resume ranker\resources\Idea Submission Template _ Redrob.pptx"

if not os.path.exists(pptx_path):
    print("PPTX not found!")
else:
    prs = Presentation(pptx_path)
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- SLIDE {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"[{shape.name}]: {shape.text}")
